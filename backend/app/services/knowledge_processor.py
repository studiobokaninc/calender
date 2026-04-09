import os
import logging
import asyncio
import shutil
import uuid
from typing import Dict, Any, List, Optional
from sqlalchemy.orm import Session
from .. import crud, models, schemas
from .llm import LLMClient
from .rag import rag_service
from .meeting_analyzer import MeetingAnalyzer
from ..database import SessionLocal
from ..timezone import now_jst_naive
from pathlib import Path

# Paths
BACKEND_ROOT = Path(__file__).resolve().parent.parent.parent
STATIC_DIR = BACKEND_ROOT / "static"
KNOWLEDGE_DIR = STATIC_DIR / "knowledge"

logger = logging.getLogger(__name__)

class KnowledgeProcessor:
    def __init__(self, api_key: str):
        self.api_key = api_key
        self.llm_client = LLMClient(api_key=api_key)
        self.meeting_analyzer = MeetingAnalyzer(api_key=api_key)

    async def process_knowledge_item(self, item_id: int):
        """Process a knowledge item based on its file type. Creating its own session for background safety."""
        db = SessionLocal()
        try:
            db_item = crud.get_knowledge_item(db, item_id=item_id)
            if not db_item:
                logger.error(f"KnowledgeItem {item_id} not found.")
                return

            crud.update_knowledge_item(db, db_item, {"status": "processing"})
            
            file_path = db_item.file_path
            # Resolve relative/virtual paths
            if file_path.startswith("/"):
                file_path = str(BACKEND_ROOT / file_path.lstrip("/"))
            elif not os.path.isabs(file_path):
                file_path = str(BACKEND_ROOT / file_path)

            content_text = ""
            summary = ""
            metadata = {}

            # Determine processor based on extension
            ext = os.path.splitext(db_item.file_name)[1].lower().lstrip(".")
            file_type = db_item.file_type 
            if ext in ["pdf"]: file_type = "pdf"
            elif ext in ["xlsx", "xls"]: file_type = "excel"
            elif ext in ["pptx", "ppt"]: file_type = "ppt"
            elif ext in ["png", "jpg", "jpeg", "webp"]: file_type = "image"
            elif ext in ["mp3", "m4a", "wav"]: file_type = "audio"
            elif ext in ["txt", "md", "json", "csv"]: file_type = "text"

            logger.info(f"Processing type: {file_type} (Ext: {ext})")

            # Prepare metadata for RAG
            # Use created_at as the reference date for knowledge items
            ref_date = db_item.created_at or now_jst_naive()
            rag_metadata = {
                "item_id": item_id, 
                "title": db_item.title, 
                "file_name": db_item.file_name,
                "project_id": db_item.project_id,
                "type": "knowledge",
                "date": ref_date.isoformat()
            }

            if file_type == "pdf":
                # Use Gemini for OCR - much better than simple local read
                print(f"KnowledgeProcessor: [{db_item.file_name}] Gemini OCR 解読を開始...")
                content_text = await self._ocr_pdf_via_gemini(file_path, is_summary=False)
                print(f"KnowledgeProcessor: [{db_item.file_name}] Gemini 要約作成を開始...")
                summary = await self._ocr_pdf_via_gemini(file_path, is_summary=True)
                
                print(f"KnowledgeProcessor: [{db_item.file_name}] RAG に PDF 内容を追加中...")
                await rag_service.add_text(content_text, metadata=rag_metadata)

            elif file_type in ["excel", "ppt"]:
                # Custom local extraction + LLM Summary
                result = await self._process_complex_doc(file_path, file_type)
                summary = result["summary"]
                content_text = result["content"]
                
                # Add to RAG (Summary + Content)
                full_kb_text = f"SUMMARY: {summary}\n\n--- FULL CONTENT ---\n{content_text}"
                print(f"KnowledgeProcessor: RAG に内容を追加中... ({db_item.file_name})")
                await rag_service.add_text(full_kb_text, metadata=rag_metadata)

            elif file_type == "image":
                content_text = await self._ocr_image(file_path)
                summary = await self._generate_summary_from_text(content_text)
                
                print(f"KnowledgeProcessor: RAG に画像テキストを追加中... ({db_item.file_name})")
                await rag_service.add_text(content_text, metadata=rag_metadata)

            elif file_type == "audio":
                res = await self.meeting_analyzer._process_segment_with_retry(file_path, 1, 1, 0)
                if res:
                    content_text = res.get("transcript", "")
                    summary = f"Summary: {res.get('summary', 'N/A')}\nDecisions: {res.get('decisions')}\nTasks: {res.get('tasks')}"
                    
                    print(f"KnowledgeProcessor: RAG にオーディオ文字起こしを追加中... ({db_item.file_name})")
                    await rag_service.add_text(content_text, metadata=rag_metadata)
            
            elif file_type == "text":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content_text = f.read()
                summary = await self._generate_summary_from_text(content_text)
                
                print(f"KnowledgeProcessor: RAG にテキストドキュメントを追加中... ({file_path})")
                await rag_service.add_text(content_text, metadata=rag_metadata)

            # Auto-tagging
            tags = await self._generate_tags(content_text or summary)
            for tag_name in tags:
                crud.add_knowledge_tag(db, item_id, tag_name)

            crud.update_knowledge_item(db, db_item, {
                "status": "completed",
                "content_text": content_text,
                "summary": summary,
                "updated_at": now_jst_naive()
            })
            db.commit()
            logger.info(f"KnowledgeItem {item_id} processed successfully.")

        except Exception as e:
            logger.error(f"Failed to process KnowledgeItem {item_id}: {e}")
            import traceback
            logger.error(traceback.format_exc())
            try:
                # Refresh db_item since it might be detached if error was DB-related
                db_item = crud.get_knowledge_item(db, item_id)
                if db_item:
                    crud.update_knowledge_item(db, db_item, {"status": "failed"})
            except:
                pass
        finally:
            db.close()

    async def _process_complex_doc(self, file_path: str, file_type: str) -> dict:
        """Local text extraction for Excel and PPTX to avoid Gemini MIME type 400 errors."""
        extracted_text = ""
        try:
            if file_type == "excel":
                import pandas as pd
                df_dict = pd.read_excel(file_path, sheet_name=None)
                for sheet_name, df in df_dict.items():
                    extracted_text += f"\nSheet: {sheet_name}\n"
                    # df.to_markdown() requires tabulate, use to_string() instead
                    extracted_text += df.to_string() + "\n"
            elif file_type == "ppt":
                from pptx import Presentation
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides):
                    extracted_text += f"\nSlide {i+1}\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            extracted_text += shape.text + "\n"
                        if shape.has_table:
                             for row in shape.table.rows:
                                extracted_text += " | ".join([cell.text for cell in row.cells]) + "\n"
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text
                        if notes: extracted_text += f"(Note: {notes})\n"
        except Exception as e:
            logger.error(f"Failed to locally extract text from {file_type}: {e}")
            extracted_text = f"Error extracting text from {file_type}."

        if not extracted_text or len(extracted_text) < 10:
             return {"summary": "Error: Could not extract content.", "content": ""}

        prompt = f"以下の{file_type}ファイルから抽出された内容を分析し、重要なポイントやデータ構造を強調した構造化された要約を【日本語で】作成してください。\n\n{extracted_text[:30000]}"
        inputs = {"mode": "utility", "no_actions": True}
        response_text = ""
        async for chunk in self.llm_client.stream_chat(prompt, f"kb_proc_{uuid.uuid4()}", inputs):
            if chunk.get("event") == "message":
                response_text += chunk.get("answer", "")
        return {"summary": response_text, "content": extracted_text}

    async def _ocr_pdf_via_gemini(self, file_path: str, is_summary: bool = False) -> str:
        prompt = "このPDFの内容を【日本語で】要約してください。" if is_summary else "このPDFのすべてのテキストをMarkdown形式で【日本語で】書き起こしてください。"
        inputs = {"mode": "admin", "attachments": [file_path], "no_actions": True}
        response_text = ""
        async for chunk in self.llm_client.stream_chat(prompt, f"kb_pdf_{uuid.uuid4()}", inputs):
            if chunk.get("event") == "message":
                response_text += chunk.get("answer", "")
        return response_text

    async def _ocr_image(self, file_path: str) -> str:
        prompt = "この画像からすべてのテキストを正確に【日本語で】書き起こしてください。表や構造がある場合は、Markdownで表現してください。"
        inputs = {"mode": "admin", "attachments": [file_path], "no_actions": True}
        response_text = ""
        async for chunk in self.llm_client.stream_chat(prompt, f"kb_ocr_{uuid.uuid4()}", inputs):
            if chunk.get("event") == "message":
                response_text += chunk.get("answer", "")
        return response_text

    async def _generate_summary_from_text(self, text: str) -> str:
        prompt = f"以下のテキストを【日本語で】簡潔に要約してください：\n\n{text}"
        inputs = {"mode": "utility", "no_actions": True}
        response_text = ""
        async for chunk in self.llm_client.stream_chat(prompt, f"kb_sum_{uuid.uuid4()}", inputs):
            if chunk.get("event") == "message":
                response_text += chunk.get("answer", "")
        return response_text

    async def _generate_summary_from_file_content(self, file_path: str, file_type: str) -> str:
        prompt = f"この{file_type}ファイルの内容を【日本語で】要約してください。"
        inputs = {"mode": "admin", "attachments": [file_path], "no_actions": True}
        response_text = ""
        async for chunk in self.llm_client.stream_chat(prompt, f"kb_sum_file_{uuid.uuid4()}", inputs):
            if chunk.get("event") == "message":
                response_text += chunk.get("answer", "")
        return response_text

    async def _generate_tags(self, text: str) -> List[str]:
        if not text: return []
        prompt = f"以下のコンテンツに関連するキーワード（タグ）を3〜5個、【日本語で】抽出してください。カンマ区切りのリスト形式でのみ出力してください。\n\n{text}"
        inputs = {"mode": "admin", "no_actions": True}
        response_text = ""
        async for chunk in self.llm_client.stream_chat(prompt, f"kb_tags_{uuid.uuid4()}", inputs):
            if chunk.get("event") == "message":
                response_text += chunk.get("answer", "")
        return [t.strip() for t in response_text.split(",") if t.strip()]
